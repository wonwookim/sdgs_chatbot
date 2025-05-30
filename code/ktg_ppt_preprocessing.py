import os
import json
import re
from pptx import Presentation

def parse_page_ranges(page_range_str):
    """
    페이지 범위 문자열을 파싱하는 함수 (ktg_preprocessing.py와 동일)
    """
    pages = set()
    ranges = page_range_str.split(',')
    for r in ranges:
        r = r.strip()
        if '-' in r:
            start, end = map(int, r.split('-'))
            pages.update(range(start, end + 1))
        else:
            pages.add(int(r))
    return pages

def get_category_metadata():
    """
    metadata_category.txt를 파싱하여 카테고리/페이지 집합 반환 (ktg_preprocessing.py와 동일)
    """
    categories = {}
    current_main = None
    all_pages = set()
    with open('metadata_category.txt', 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            if line[0].isdigit():
                main_cat = line.split('.')[1].strip()
                current_main = main_cat
                categories[current_main] = {'subcategories': {}}
            elif line[0].isalpha() and current_main:
                sub_cat = line.split('.')[1].strip()
                page_range = re.search(r'\(([0-9,\- ]+)\)', line)
                if page_range:
                    pages = parse_page_ranges(page_range.group(1))
                    categories[current_main]['subcategories'][sub_cat] = list(pages)
                    all_pages.update(pages)
    return categories, all_pages

def clean_text(text):
    """
    텍스트를 정제하는 함수 (PDF용과 동일하게 적용)
    Args:
        text (str): 정제할 텍스트
    Returns:
        str: 정제된 텍스트
    """
    # 연속된 공백(단, 줄바꿈은 보존) 제거
    text = re.sub(r'[ \t\x0b\x0c\r]+', ' ', text)
    # ●, • 등 특정 특수문자 제거
    text = re.sub(r'[●•]', '', text)
    # 특수문자 제거 (단, 한글, 영문, 숫자, 괄호, 줄바꿈, 기본 문장부호는 유지)
    text = re.sub(r'[^\w\s가-힣ㄱ-ㅎㅏ-ㅣ.,!?;:()\[\]{}""\'\'\-\n]', '', text)
    # 연속된 줄바꿈 정리
    text = re.sub(r'\n{3,}', '\n\n', text)
    # 문단 시작 부분의 불필요한 공백 제거
    text = re.sub(r'\n\s+', '\n', text)
    # 문장 끝 부분의 불필요한 공백 제거
    text = re.sub(r'\s+\n', '\n', text)
    # 앞뒤 공백 제거
    text = text.strip()
    return text

def extract_text_from_pptx(pptx_path):
    """
    PPTX 파일에서 슬라이드별 텍스트 추출 및 정제 + 메타데이터(우측 상단 페이지번호, 카테고리)
    """
    prs = Presentation(pptx_path)
    slides_data = []
    categories, all_pages = get_category_metadata()
    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        page_number = None
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        # 우측 상단 텍스트 상자에서 숫자 추출
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                # shape 위치 정보
                left = shape.left
                top = shape.top
                # 우측 상단 기준: left >= 70% slide_width, top <= 20% slide_height
                if left >= int(slide_width * 0.7) and top <= int(slide_height * 0.2):
                    nums = re.findall(r'\b\d{1,3}\b', shape.text)
                    if nums:
                        page_number = int(nums[0])
                        break
        # 전체 텍스트 추출
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        raw_text = "\n".join(texts)
        cleaned_text = clean_text(raw_text)
        # page_number가 없으면 전체 텍스트에서 추출(백업)
        if page_number is None:
            page_numbers = re.findall(r'\b\d{1,3}\b', cleaned_text)
            page_number = int(page_numbers[0]) if page_numbers else idx
        # 카테고리 할당
        page_categories = []
        for main_cat, main_data in categories.items():
            for sub_cat, pages in main_data['subcategories'].items():
                if page_number in pages:
                    page_categories.append({
                        'main_category': main_cat,
                        'sub_category': sub_cat
                    })
        slide_data = {
            "slide_number": idx,
            "page_number": page_number,
            "content": cleaned_text,
            "metadata": {
                "source": os.path.basename(pptx_path),
                "slide": idx,
                "page": page_number,
                "total_slides": len(prs.slides),
                "categories": page_categories
            }
        }
        slides_data.append(slide_data)
    return slides_data

def save_text_to_file(slides_data, output_path):
    """
    텍스트를 파일로 저장하는 함수 (json, txt)
    Args:
        slides_data (list): 슬라이드별 텍스트와 메타데이터가 포함된 딕셔너리 리스트
        output_path (str): 저장할 파일 경로
    """
    # JSON 형식으로 저장
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(slides_data, f, ensure_ascii=False, indent=2)
    # 텍스트 파일로도 저장 (메타데이터 포함)
    txt_path = output_path.replace('.json', '.txt')
    with open(txt_path, 'w', encoding='utf-8') as f:
        for slide in slides_data:
            f.write(f"\n=== Slide {slide['slide_number']} (Page {slide['page_number']}) ===\n")
            # 카테고리 정보 추가
            if slide['metadata']['categories']:
                f.write("Categories:\n")
                for cat in slide['metadata']['categories']:
                    f.write(f"- {cat['main_category']} > {cat['sub_category']}\n")
            f.write(slide['content'])
            f.write('\n\n')

def main():
    pptx_path = "KT&G_지속가능경영보고서_편집본(2023).pptx"
    output_path = "ktg_ppt_report_2023.json"
    print("PPTX에서 슬라이드별 텍스트 추출 및 정제 중...")
    slides_data = extract_text_from_pptx(pptx_path)
    print(f"총 {len(slides_data)}개 슬라이드 처리 완료. 파일 저장 중...")
    save_text_to_file(slides_data, output_path)
    print(f"완료! 결과가 {output_path}와 {output_path.replace('.json', '.txt')}에 저장되었습니다.")

if __name__ == "__main__":
    main() 