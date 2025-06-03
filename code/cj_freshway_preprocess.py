import os
import json
import re
import unicodedata
from pptx import Presentation

def parse_page_ranges(page_range_str):
    """
    페이지 범위 문자열을 파싱하는 함수 (예: '3-9,11,13-15' -> [3,4,5,6,7,8,9,11,13,14,15])
    """
    pages = set()
    for part in page_range_str.split(','):
        part = part.strip()
        if '-' in part:
            start, end = map(int, part.split('-'))
            pages.update(range(start, end+1))
        elif part:
            pages.add(int(part))
    return sorted(pages)

def get_category_metadata(metadata_path):
    """
    metadata_category.txt를 파싱하여 카테고리/페이지 집합 반환
    파일 예시: main_category|sub_category|page_range
    """
    category_map = {}
    with open(metadata_path, encoding='utf-8') as f:
        for line in f:
            if not line.strip() or line.startswith('#'):
                continue
            main_cat, sub_cat, page_range = line.strip().split('|')
            for page in parse_page_ranges(page_range):
                category_map[page] = {
                    "main_category": main_cat,
                    "sub_category": f"{sub_cat}({page_range})"
                }
    return category_map

def clean_text(text):
    """
    텍스트를 정제하는 함수 (PDF용과 동일하게 적용)
    """
    text = unicodedata.normalize('NFKC', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def extract_text_from_pptx(pptx_path, category_map):
    """
    PPTX 파일에서 슬라이드별 텍스트 추출 및 정제 + 메타데이터(카테고리 등)
    content는 shape별 줄바꿈을 살려서 보기 좋게 만듦
    pdf_page_number는 content의 첫 번째 숫자만 있는 줄에서 추출
    """
    prs = Presentation(pptx_path)
    slides_data = []
    total_slides = len(prs.slides)
    source_name = os.path.basename(pptx_path)
    for idx, slide in enumerate(prs.slides):
        slide_number = idx + 1
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for line in shape.text.splitlines():
                    clean = clean_text(line)
                    if clean:
                        lines.append(clean)
        # pdf_page_number 추출: 첫 번째로 등장하는 숫자만 있는 줄
        pdf_page_number = slide_number
        pdf_page_number_found = None
        new_lines = []
        for line in lines:
            if pdf_page_number_found is None and line.isdigit():
                pdf_page_number_found = int(line)
                continue  # pdf_page_number 줄은 content에서 제외
            new_lines.append(line)
        if pdf_page_number_found is not None:
            pdf_page_number = pdf_page_number_found
        content = "\n".join(new_lines)
        cat = category_map.get(slide_number, {})
        slide_data = {
            "slide_number": slide_number,
            "pdf_page_number": pdf_page_number,
            "content": content,
            "metadata": {
                "source": source_name,
                "slide": slide_number,
                "pdf_page_number": pdf_page_number,
                "total_slides": total_slides,
                "categories": [cat] if cat else []
            }
        }
        slides_data.append(slide_data)
    return slides_data

def save_text_to_file(slides_data, output_path):
    """
    텍스트를 파일로 저장하는 함수 (json)
    """
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(slides_data, f, ensure_ascii=False, indent=2)
    print(f"저장 완료: {output_path}")

def save_text_to_txt(slides_data, output_path):
    """
    슬라이드별 텍스트를 하나의 txt 파일로 저장
    """
    with open(output_path, "w", encoding="utf-8") as f:
        for slide in slides_data:
            f.write(f"--- Slide {slide['slide_number']} (PDF Page {slide['pdf_page_number']}) ---\n")
            f.write(slide['content'])
            f.write("\n\n")
    print(f"텍스트 파일 저장 완료: {output_path}")

def process_pptx_directory(input_dir, metadata_path, output_path):
    category_map = get_category_metadata(metadata_path)
    all_data = []
    for fname in os.listdir(input_dir):
        if fname.lower().endswith('.pptx'):
            pptx_path = os.path.join(input_dir, fname)
            print(f"Processing {pptx_path}")
            slides = extract_text_from_pptx(pptx_path, category_map)
            all_data.extend(slides)
    save_text_to_file(all_data, output_path)
    # 텍스트 파일로도 저장
    txt_output_path = output_path.replace('.json', '.txt')
    save_text_to_txt(all_data, txt_output_path)

if __name__ == "__main__":
    input_dir = os.path.join("sdgs_chatbot", "code", "data", "ppt")
    metadata_path = os.path.join("sdgs_chatbot", "code", "data", "metadata_category.txt")
    output_path = os.path.join("sdgs_chatbot", "code", "data", "processed", "cj_freshway.json")
    process_pptx_directory(input_dir, metadata_path, output_path)
