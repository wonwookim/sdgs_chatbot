from pptx import Presentation
import re
from pathlib import Path
import json
import os
import chromadb
from chromadb.utils import embedding_functions
import uuid
from langchain.text_splitter import RecursiveCharacterTextSplitter

def get_section_and_subsection(page_num, ppt_name):
    """페이지 번호와 PPT 이름에 따라 섹션과 서브섹션을 결정하는 함수
       각 섹션과 서브섹션은 ppt 파일에 따라 다르게 정의되어 있음
       
       예시: 
        - cj_section_data.txt
        
            [ESG 전략]
            1-2:ESG 전략

            [Environment]
            3-9:기후변화 대응
            10-13:지속 가능한 유통
            14-16:환경 경영
            
        - shinhan_section_data.txt
        
            [Environment]
            1:Environmental
            2:환경경영 추진 체계
            3-4:내부 탄소배출량 관리
            5-6:금융 배출량 관리
            7-11:친환경 금융 확대
            12-13:친환경 실천 교육 및 캠페인

            [Social]
            14:Social
            15-27:혁신 및 포용금융
            28-32:지역사회 발전 및 투자
            33-37:고객 만족도 제고
            38-45:인재 경영
            46-48:인권 및 다양성
            49-51:안전 및 보건
            52:지속가능한 공급망
    """
    # PPT 파일명에서 확장자를 제외한 이름을 추출
    base_name = Path(ppt_name).stem
    section_file = f'data/config/{base_name}_section_data.txt'
    
    try: # 섹션 데이터 파일 읽기
        with open(section_file, 'r', encoding='utf-8') as f:
            section_data = f.read()
    except FileNotFoundError: # 섹션 데이터 파일이 없으면 'Other' 반환
        print(f"Warning: Section data file {section_file} not found.")
        return 'Other', 'Other'

    current_section = 'Other'
    current_subsection = 'Other'

    # 섹션 데이터 파싱
    sections = section_data.split('\n\n')
    for section in sections:
        if not section.strip(): # 섹션이 비어있을 경우
            continue
        
        lines = section.strip().split('\n')
        section_name = lines[0].strip('[]')
        
        for line in lines[1:]:
            if not line.strip():
                continue
                
            page_range, subsection_name = line.split(':')
            
            # 페이지 범위 처리
            if '-' in page_range:
                start, end = map(int, page_range.split('-'))
                if start <= page_num <= end:
                    current_section = section_name
                    current_subsection = subsection_name
                    break
            else:
                if int(page_range) == page_num:
                    current_section = section_name
                    current_subsection = subsection_name
                    break
    # 현재 보고있는 페이지의 section과 subsection을 반환
    return current_section, current_subsection

def preprocess_text(text):
    """텍스트 전처리 함수"""
    # 여러 줄의 공백을 하나의 줄바꿈으로 변경
    text = re.sub(r'\n\s*\n', '\n', text)
    
    # 불필요한 공백 제거
    text = re.sub(r'\s+', ' ', text)
    
    # 특수문자 처리 (예: 불필요한 특수문자 제거, 필요한 것은 유지)
    text = re.sub(r'[^\w\s\.,!?\/\-\(\)\[\]\{\}:;\'\"]+', '', text)
    
    # 중복된 문장부호 제거
    text = re.sub(r'[.]{2,}', '.', text)
    text = re.sub(r'[!]{2,}', '!', text)
    text = re.sub(r'[?]{2,}', '?', text)
    
    # 앞뒤 공백 제거
    text = text.strip()
    
    # 빈 줄로만 이루어진 경우 None 반환
    if not text or text.isspace():
        return None
        
    return text

def extract_slide_text(slide):
    """슬라이드에서 텍스트를 추출하는 함수"""
    text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            # 텍스트 전처리 수행
            processed_text = preprocess_text(text)
            if processed_text:
                text_parts.append(processed_text)
    
    # 모든 텍스트를 합치고 다시 한번 전처리
    combined_text = '\n'.join(text_parts)
    return preprocess_text(combined_text) if combined_text else None

def process_ppt(ppt_path):
    """PPT 파일을 처리하고 청크를 생성하는 메인 함수
       문서 특성 상 섹션과 서브섹션에 따라 문맥이 유지될 수 있기 때문에, 섹션 & 서브섹션으로 chunck를 구성
    """
    print(f"PPT 파일을 읽는 중: {ppt_path}")
    presentation = Presentation(ppt_path) # PPT 파일 읽기
    total_slides = len(presentation.slides) # PPT 총 슬라이드 수
    print(f"PPT 총 슬라이드 수: {total_slides}")
    
    # PPT 파일명을 소스 이름으로 사용
    ppt_name = Path(ppt_path).name
    source_name = Path(ppt_path).stem
    
    # 섹션과 서브섹션별로 텍스트를 그룹화하는 딕셔너리
    section_texts = {}
    # 각 텍스트의 페이지 번호를 추적하는 딕셔너리
    text_page_numbers = {}
    
    # 먼저 모든 텍스트를 섹션과 서브섹션별로 수집
    for idx, slide in enumerate(presentation.slides):
        page_num = idx + 1
        print(f"슬라이드 {page_num} 처리 중...")
        
        section_type, sub_section = get_section_and_subsection(page_num, ppt_name)
        print(f"섹션: {section_type}, 서브섹션: {sub_section}")
        
        slide_text = extract_slide_text(slide)
        
        if slide_text: # 텍스트가 있으면
            section_key = (section_type, sub_section)
            if section_key not in section_texts:
                section_texts[section_key] = []
                text_page_numbers[section_key] = []
            
            section_texts[section_key].append(slide_text)
            text_page_numbers[section_key].append(page_num)
    
    all_chunks = []
    
    # 각 섹션/서브섹션 그룹에 대해 처리
    for (section_type, sub_section), texts in section_texts.items():
        # 해당 섹션의 모든 텍스트를 결합
        combined_text = "\n".join(texts)
        
        # RecursiveCharacterTextSplitter로 분할
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=50,
            length_function=len,
            separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
        )
        
        split_texts = text_splitter.split_text(combined_text)
        
        # 각 청크에 대한 메타데이터 생성
        page_numbers = text_page_numbers[(section_type, sub_section)]
        page_range = f"{min(page_numbers)}-{max(page_numbers)}"  # 페이지 범위를 문자열로 변환
        
        for i, chunk_text in enumerate(split_texts):
            chunk = {
                'text': chunk_text,
                'metadata': {
                    'section': section_type, # 섹션
                    'sub_section': sub_section, # 서브섹션
                    'source': source_name.upper(), # 회사 이름
                    'page_range': page_range,  # 페이지 범위를 문자열로 저장
                    'chunk_index': i, # 청크 인덱스
                    'total_chunks_in_section': len(split_texts) # 섹션 내 청크 수
                }
            }
            all_chunks.append(chunk)
            
        print(f"{section_type} 섹션({sub_section})에서 {len(split_texts)}개의 청크 생성됨")
    
    return all_chunks, source_name

def save_to_chroma(chunks, collection_name):
    """청크를 ChromaDB에 저장하는 함수"""
    # ChromaDB 클라이언트 초기화
    client = chromadb.PersistentClient(path="./data/chromadb")
    
    # openai의 embedding 함수는 비용 발생하기 때문에, 비용 발생하지 않는 함수 사용
    embedding_function = embedding_functions.SentenceTransformerEmbeddingFunction(
        model_name="jhgan/ko-sroberta-multitask"
    )
    
    # 컬렉션이 이미 존재하면 삭제
    try:
        client.delete_collection(name=collection_name)
    except:
        pass
    
    # 새로운 컬렉션 생성
    collection = client.create_collection(
        name=collection_name,
        embedding_function=embedding_function
    )
    
    # 청크 데이터 준비
    ids = []
    texts = []
    metadatas = []
    
    for chunk in chunks:
        chunk_id = str(uuid.uuid4())
        ids.append(chunk_id)
        texts.append(chunk['text'])
        metadatas.append(chunk['metadata'])
    
    # 데이터 추가
    collection.add(
        ids=ids,
        documents=texts,
        metadatas=metadatas
    )
    
    print(f"ChromaDB에 {len(chunks)}개의 청크가 저장되었습니다.")
    return collection

def main():
    ppt_dir = "data/ppts"
    print("프로그램 시작")
    
    all_ppt_chunks = []
    
    # data/ppts 디렉토리의 모든 PPT/PPTX 파일 처리
    for ppt_file in Path(ppt_dir).glob("*.ppt*"):
        print(f"\n{ppt_file.name} 처리 시작")
        chunks, source_name = process_ppt(str(ppt_file))
        
        if chunks:
            # JSON 파일로 저장
            output_path = f"outputs/{source_name}_chunk.json"
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(chunks, f, ensure_ascii=False, indent=2)
            print(f"데이터가 {len(chunks)}개의 청크로 나뉘어 {output_path}에 저장되었습니다.")
            
            # 모든 청크를 리스트에 추가
            all_ppt_chunks.extend(chunks)
        else:
            print(f"{ppt_file.name}에서 처리할 수 있는 텍스트를 찾을 수 없습니다.")
    
    # 모든 PPT의 청크를 하나의 ChromaDB 컬렉션에 저장
    if all_ppt_chunks:
        collection_name = "ppt_documents_collection"
        collection = save_to_chroma(all_ppt_chunks, collection_name)
        print(f"총 {len(all_ppt_chunks)}개의 청크가 통합 컬렉션에 저장되었습니다.")

if __name__ == "__main__":
    main() 