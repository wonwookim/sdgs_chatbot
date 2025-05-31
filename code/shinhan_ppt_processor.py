from pptx import Presentation
import re
from pathlib import Path
import json
import os
import chromadb
from chromadb.utils import embedding_functions
import uuid

def get_section_and_subsection(page_num):
    """페이지 번호에 따라 섹션과 서브섹션을 결정하는 함수"""
    # Environmental (p.1 ~ 13)
    if 1 <= page_num <= 13:
        section = 'Environmental'
        if page_num == 2:
            subsection = '환경경영 추진 체계'
        elif 3 <= page_num <= 4:
            subsection = '내부 탄소배출량 관리'
        elif 5 <= page_num <= 6:
            subsection = '금융 배출량 관리'
        elif 7 <= page_num <= 11:
            subsection = '친환경 금융 확대'
        elif 12 <= page_num <= 13:
            subsection = '친환경 실천 교육 및 캠페인'
        else:
            subsection = section  # None 대신 섹션 값 사용
    
    # Social (p.14 ~ 52)
    elif 14 <= page_num <= 52:
        section = 'Social'
        if 15 <= page_num <= 27:
            subsection = '혁신 및 포용금융'
        elif 28 <= page_num <= 32:
            subsection = '지역사회 발전 및 투자'
        elif 33 <= page_num <= 37:
            subsection = '고객 만족도 제고'
        elif 38 <= page_num <= 45:
            subsection = '인재 경영'
        elif 46 <= page_num <= 48:
            subsection = '인권 및 다양성'
        elif 49 <= page_num <= 51:
            subsection = '안전 및 보건'
        elif page_num == 52:
            subsection = '지속가능한 공급망'
        else:
            subsection = section  # None 대신 섹션 값 사용
    
    # Governance (p.53 ~ 69)
    elif 53 <= page_num <= 69:
        section = 'Governance'
        if 54 <= page_num <= 58:
            subsection = '건전한 지배구조'
        elif 59 <= page_num <= 61:
            subsection = '윤리 및 준법 경영'
        elif 62 <= page_num <= 64:
            subsection = '소비자 권익 보호'
        elif 65 <= page_num <= 67:
            subsection = '개인정보보호 및 데이터 보안'
        elif 68 <= page_num <= 69:
            subsection = '리스크 관리'
        else:
            subsection = section  # None 대신 섹션 값 사용
    
    else:
        section = 'Other'
        subsection = section  # None 대신 섹션 값 사용
    
    return section, subsection

def preprocess_text(text):
    """텍스트 전처리 함수"""
    # 여러 줄의 공백을 하나의 줄바꿈으로 변경
    text = re.sub(r'\n\s*\n', '\n', text)
    
    # 불필요한 공백 제거
    text = re.sub(r'\s+', ' ', text)
    
    # 특수문자 처리 (예: 불필요한 특수문자 제거, 필요한 것은 유지)
    text = re.sub(r'[^\w\s\.,!?\/\-\(\)\[\]\{\}:;\'\"]+', '', text)
    
    # 중복된 문장부호 제거
    text = re.sub(r'[.]{2,}', '.', text)
    text = re.sub(r'[!]{2,}', '!', text)
    text = re.sub(r'[?]{2,}', '?', text)
    
    # 앞뒤 공백 제거
    text = text.strip()
    
    # 빈 줄로만 이루어진 경우 None 반환
    if not text or text.isspace():
        return None
        
    return text

def extract_slide_text(slide):
    """슬라이드에서 텍스트를 추출하는 함수"""
    text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            # 텍스트 전처리 수행
            processed_text = preprocess_text(text)
            if processed_text:
                text_parts.append(processed_text)
    
    # 모든 텍스트를 합치고 다시 한번 전처리
    combined_text = '\n'.join(text_parts)
    return preprocess_text(combined_text) if combined_text else None

def create_chunks(text, section_type, sub_section, page_num, chunk_size=1000):
    """텍스트를 청크로 나누는 함수"""
    # 문장 단위로 분리 (마침표, 느낌표, 물음표 다음에 공백이 오는 경우)
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_length = 0
    
    for sentence in sentences:
        # 각 문장도 전처리
        sentence = preprocess_text(sentence)
        if not sentence:
            continue
            
        if current_length + len(sentence) > chunk_size:
            if current_chunk:
                chunk_text = ' '.join(current_chunk)
                # 청크 전체에 대해서도 한번 더 전처리
                chunk_text = preprocess_text(chunk_text)
                if chunk_text:
                    chunk = {
                        'text': chunk_text,
                        'metadata': {
                            'section': section_type,
                            'sub_section': sub_section,
                            'source': '신한라이프 2023 ESG 보고서',
                            'page': page_num
                        }
                    }
                    chunks.append(chunk)
            current_chunk = [sentence]
            current_length = len(sentence)
        else:
            current_chunk.append(sentence)
            current_length += len(sentence)
    
    if current_chunk:
        chunk_text = ' '.join(current_chunk)
        # 마지막 청크도 전처리
        chunk_text = preprocess_text(chunk_text)
        if chunk_text:
            chunk = {
                'text': chunk_text,
                'metadata': {
                    'section': section_type,
                    'sub_section': sub_section,
                    'source': '신한라이프 2023 ESG 보고서',
                    'page': page_num
                }
            }
            chunks.append(chunk)
    
    return chunks

def save_to_chroma(chunks):
    """청크를 ChromaDB에 저장하는 함수"""
    # ChromaDB 클라이언트 초기화
    client = chromadb.PersistentClient(path="./outputs/chromadb")
    
    # OpenAI의 임베딩 함수 사용
    embedding_function = embedding_functions.SentenceTransformerEmbeddingFunction(
        model_name="jhgan/ko-sroberta-multitask"
    )
    
    # 컬렉션 생성 또는 가져오기
    collection_name = "shinhan_life_esg_2023"
    try:
        collection = client.get_collection(
            name=collection_name,
            embedding_function=embedding_function
        )
        # 기존 데이터 삭제
        collection.delete(where={})
    except:
        collection = client.create_collection(
            name=collection_name,
            embedding_function=embedding_function
        )
    
    # 청크 데이터 준비
    ids = []
    texts = []
    metadatas = []
    
    for chunk in chunks:
        chunk_id = str(uuid.uuid4())
        ids.append(chunk_id)
        texts.append(chunk['text'])
        metadatas.append(chunk['metadata'])
    
    # 데이터 추가
    collection.add(
        ids=ids,
        documents=texts,
        metadatas=metadatas
    )
    
    print(f"ChromaDB에 {len(chunks)}개의 청크가 저장되었습니다.")
    return collection

def process_ppt(ppt_path):
    """PPT 파일을 처리하고 청크를 생성하는 메인 함수"""
    print(f"PPT 파일을 읽는 중: {ppt_path}")
    presentation = Presentation(ppt_path)
    total_slides = len(presentation.slides)
    print(f"PPT 총 슬라이드 수: {total_slides}")
    
    all_chunks = []
    
    for idx, slide in enumerate(presentation.slides):
        page_num = idx + 1
        print(f"슬라이드 {page_num} 처리 중...")
        
        section_type, sub_section = get_section_and_subsection(page_num)
        print(f"섹션: {section_type}, 서브섹션: {sub_section}")
        
        slide_text = extract_slide_text(slide)
        
        if slide_text:  # None 체크
            chunks = create_chunks(slide_text, section_type, sub_section, page_num)
            if chunks:  # 빈 리스트 체크
                all_chunks.extend(chunks)
                print(f"{section_type} 섹션({sub_section})에서 {len(chunks)}개의 청크 생성됨")
    
    return all_chunks

def main():
    ppt_path = "data/ppts/2023 신한라이프 지속가능경영보고서.pptx"
    print("프로그램 시작")
    
    chunks = process_ppt(ppt_path)
    
    if chunks:
        # JSON 파일로 저장
        output_path = "outputs/shinhan_chunk.json"
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(chunks, f, ensure_ascii=False, indent=2)
        print(f"ESG 데이터가 {len(chunks)}개의 청크로 나뉘어 {output_path}에 저장되었습니다.")
        
        # ChromaDB에 저장
        collection = save_to_chroma(chunks)
        
    else:
        print("처리할 수 있는 텍스트를 찾을 수 없습니다.")

if __name__ == "__main__":
    main() 