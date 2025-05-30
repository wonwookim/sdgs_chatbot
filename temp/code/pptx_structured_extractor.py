import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def get_shape_text_info(shape):
    """Extract text, color, bold, and font size from a shape (if available)."""
    if not hasattr(shape, "text") or not shape.text.strip():
        return None
    text = shape.text.strip()
    color = None
    bold = None
    size = None
    # Try to get font info from the first run of the first paragraph
    try:
        para = shape.text_frame.paragraphs[0]
        if para.runs:
            run = para.runs[0]
            font = run.font
            if font.color and font.color.rgb:
                color = str(font.color.rgb)
            bold = font.bold
            if font.size:
                size = font.size.pt
    except Exception:
        pass
    # Try to get fill color (background)
    fill_color = None
    try:
        if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
            fill_color = str(shape.fill.fore_color.rgb)
    except Exception:
        pass
    return {
        "text": text,
        "color": color,
        "bold": bold,
        "size": size,
        "fill_color": fill_color
    }

def classify_shape(info):
    """Classify the shape based on font size, bold, and color."""
    if not info:
        return None
    # Section: 큰 글씨, bold, 또는 fill_color가 없음
    if info["size"] and info["size"] >= 24:
        return "section"
    # Group: 색이 있거나(초록 등), fill_color가 있으면
    if info["fill_color"] and info["fill_color"] != "None":
        return "group"
    # Subsection: bold
    if info["bold"]:
        return "subsection"
    # 일반 텍스트
    return "text"

def extract_slide_structured(slide):
    slide_struct = []
    for shape in slide.shapes:
        info = get_shape_text_info(shape)
        if not info:
            continue
        shape_type = classify_shape(info)
        entry = {
            "type": shape_type,
            "text": info["text"]
        }
        if info["color"]:
            entry["color"] = info["color"]
        if info["bold"] is not None:
            entry["bold"] = info["bold"]
        if info["size"]:
            entry["size"] = info["size"]
        if info["fill_color"]:
            entry["fill_color"] = info["fill_color"]
        slide_struct.append(entry)
    return slide_struct

def extract_pptx_structured(pptx_path):
    prs = Presentation(pptx_path)
    all_slides = []
    for idx, slide in enumerate(prs.slides):
        slide_content = extract_slide_structured(slide)
        if slide_content:
            all_slides.append({
                "slide_number": idx + 1,
                "contents": slide_content
            })
    return all_slides

if __name__ == "__main__":
    pptx_path = "temp/data/ppts/2023 신한라이프 지속가능경영보고서.pptx"
    slides = extract_pptx_structured(pptx_path)
    with open("temp/data/processed/pptx_structured_meta.json", "w", encoding="utf-8") as f:
        json.dump(slides, f, ensure_ascii=False, indent=2)
    print("구조화된 메타데이터가 pptx_structured_meta.json에 저장되었습니다.") 