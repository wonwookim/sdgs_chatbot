import json
from pptx import Presentation
import re

class StructuredElement:
    def __init__(self, type, content, bbox, color=None, bold=None, size=None, fill_color=None, children=None):
        self.type = type
        self.content = content
        self.bbox = bbox
        self.color = color
        self.bold = bold
        self.size = size
        self.fill_color = fill_color
        self.children = children if children is not None else []

    def to_dict(self):
        d = {
            "type": self.type,
            "content": self.content,
            "bbox": self.bbox,
            "children": [c.to_dict() for c in self.children]
        }
        if self.color: d["color"] = self.color
        if self.bold is not None: d["bold"] = self.bold
        if self.size: d["size"] = self.size
        if self.fill_color: d["fill_color"] = self.fill_color
        return d

def get_shape_info(shape):
    text = shape.text.strip() if hasattr(shape, "text") and shape.text else ""
    bbox = [int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height)]
    color = None
    bold = None
    size = None
    fill_color = None
    try:
        para = shape.text_frame.paragraphs[0]
        if para.runs:
            run = para.runs[0]
            font = run.font
            if font.color and font.color.rgb:
                color = str(font.color.rgb)
            bold = font.bold
            if font.size:
                size = font.size.pt
    except Exception:
        pass
    try:
        if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
            fill_color = str(shape.fill.fore_color.rgb)
    except Exception:
        pass
    return text, bbox, color, bold, size, fill_color

def classify_shape(type, size, bold, fill_color):
    if size and size >= 24:
        return "section"
    if fill_color and fill_color != "None":
        return "group"
    if bold:
        return "subsection"
    return "text"

def extract_structured_elements(slide):
    elements = []
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        text, bbox, color, bold, size, fill_color = get_shape_info(shape)
        elem_type = classify_shape(shape.shape_type, size, bold, fill_color)
        elem = StructuredElement(
            type=elem_type,
            content=text,
            bbox=bbox,
            color=color,
            bold=bold,
            size=size,
            fill_color=fill_color
        )
        elements.append(elem)
    return elements

def create_chunks(text, section, slide_number, structured_elements, chunk_size=1000):
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_length = 0
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        if current_length + len(sentence) > chunk_size:
            if current_chunk:
                chunk_text = ' '.join(current_chunk)
                chunks.append({
                    'text': chunk_text,
                    'metadata': {
                        'section': section,
                        'slide_number': slide_number,
                        'source': '신한라이프 2023 ESG 보고서',
                        'structured_elements': [e.to_dict() for e in structured_elements]
                    }
                })
            current_chunk = [sentence]
            current_length = len(sentence)
        else:
            current_chunk.append(sentence)
            current_length += len(sentence)
    if current_chunk:
        chunk_text = ' '.join(current_chunk)
        chunks.append({
            'text': chunk_text,
            'metadata': {
                'section': section,
                'slide_number': slide_number,
                'source': '신한라이프 2023 ESG 보고서',
                'structured_elements': [e.to_dict() for e in structured_elements]
            }
        })
    return chunks

def extract_pptx_chunks_by_section_with_structure(pptx_path):
    prs = Presentation(pptx_path)
    section_ranges = {
        'Environmental': (24, 38),
        'Social': (40, 79),
        'Governance': (81, 99),
    }
    all_chunks = []
    for section, (start, end) in section_ranges.items():
        for idx in range(start, end + 1):
            if idx < len(prs.slides):
                slide = prs.slides[idx]
                # 슬라이드의 모든 텍스트 합치기
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                joined_text = '\n'.join(slide_text)
                # 계층 구조 추출
                structured_elements = extract_structured_elements(slide)
                if joined_text.strip():
                    chunks = create_chunks(joined_text, section, idx + 1, structured_elements)
                    all_chunks.extend(chunks)
    return all_chunks

if __name__ == "__main__":
    pptx_path = "temp/data/ppts/2023 신한라이프 지속가능경영보고서.pptx"
    chunks = extract_pptx_chunks_by_section_with_structure(pptx_path)
    with open("temp/data/processed/pptx_chunks_with_structure.json", "w", encoding="utf-8") as f:
        json.dump(chunks, f, ensure_ascii=False, indent=2)
    print("PPTX 청크+구조화 메타데이터가 pptx_chunks_with_structure.json에 저장되었습니다.") 