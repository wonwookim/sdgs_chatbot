import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

class StructuredElement:
    def __init__(self, type, content, bbox, color=None, bold=None, size=None, fill_color=None, children=None):
        self.type = type
        self.content = content
        self.bbox = bbox
        self.color = color
        self.bold = bold
        self.size = size
        self.fill_color = fill_color
        self.children = children if children is not None else []

    def to_dict(self):
        d = {
            "type": self.type,
            "content": self.content,
            "bbox": self.bbox,
            "children": [c.to_dict() for c in self.children]
        }
        if self.color: d["color"] = self.color
        if self.bold is not None: d["bold"] = self.bold
        if self.size: d["size"] = self.size
        if self.fill_color: d["fill_color"] = self.fill_color
        return d

def get_shape_info(shape):
    text = shape.text.strip() if hasattr(shape, "text") and shape.text else ""
    bbox = [int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height)]
    color = None
    bold = None
    size = None
    fill_color = None
    try:
        para = shape.text_frame.paragraphs[0]
        if para.runs:
            run = para.runs[0]
            font = run.font
            if font.color and font.color.rgb:
                color = str(font.color.rgb)
            bold = font.bold
            if font.size:
                size = font.size.pt
    except Exception:
        pass
    try:
        if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
            fill_color = str(shape.fill.fore_color.rgb)
    except Exception:
        pass
    return text, bbox, color, bold, size, fill_color

def classify_shape(type, size, bold, fill_color):
    if size and size >= 24:
        return "section"
    if fill_color and fill_color != "None":
        return "group"
    if bold:
        return "subsection"
    return "text"

def extract_structured_elements(slide):
    elements = []
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        text, bbox, color, bold, size, fill_color = get_shape_info(shape)
        elem_type = classify_shape(shape.shape_type, size, bold, fill_color)
        elem = StructuredElement(
            type=elem_type,
            content=text,
            bbox=bbox,
            color=color,
            bold=bold,
            size=size,
            fill_color=fill_color
        )
        elements.append(elem)
    return elements

def extract_pptx_structured_meta_by_section(pptx_path):
    prs = Presentation(pptx_path)
    section_ranges = {
        'Environmental': (24, 38),  # 25~39
        'Social': (40, 79),         # 41~80
        'Governance': (81, 99),     # 82~100
    }
    section_data = {k: [] for k in section_ranges}
    for section, (start, end) in section_ranges.items():
        for idx in range(start, end + 1):
            if idx < len(prs.slides):
                slide = prs.slides[idx]
                elements = extract_structured_elements(slide)
                section_data[section].append({
                    "slide_number": idx + 1,
                    "elements": [e.to_dict() for e in elements]
                })
    return section_data

if __name__ == "__main__":
    pptx_path = "temp/data/ppts/2023 신한라이프 지속가능경영보고서.pptx"
    sectioned = extract_pptx_structured_meta_by_section(pptx_path)
    with open("temp/data/processed/pptx_structured_elements_by_section.json", "w", encoding="utf-8") as f:
        json.dump(sectioned, f, ensure_ascii=False, indent=2)
    print("섹션별 구조화된 PPTX 메타데이터가 pptx_structured_elements_by_section.json에 저장되었습니다.") 