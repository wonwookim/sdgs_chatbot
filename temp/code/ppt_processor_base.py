import json
from pptx import Presentation
import re
from pathlib import Path

def extract_slide_content(slide):
    """슬라이드의 텍스트 내용을 추출합니다."""
    text_content = []
    
    # 슬라이드 제목 추출
    if slide.shapes.title:
        text_content.append(slide.shapes.title.text)
    
    # 슬라이드의 모든 텍스트 추출
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            if shape != slide.shapes.title:  # 제목은 이미 추가했으므로 제외
                text_content.append(shape.text)
    
    return "\n".join(text_content)

def create_chunks_base(text, section_type, sub_section, chunk_size=1000):
    """텍스트를 청크로 나눕니다."""
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_length = 0
    
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
            
        if current_length + len(sentence) > chunk_size:
            if current_chunk:
                chunk_text = ' '.join(current_chunk)
                chunk = {
                    'text': chunk_text,
                    'metadata': {
                        'section': section_type,
                        'sub_section': sub_section,
                        'source': '신한라이프 2023 ESG 보고서'
                    }
                }
                chunks.append(chunk)
            current_chunk = [sentence]
            current_length = len(sentence)
        else:
            current_chunk.append(sentence)
            current_length += len(sentence)
    
    if current_chunk:
        chunk_text = ' '.join(current_chunk)
        chunk = {
            'text': chunk_text,
            'metadata': {
                'section': section_type,
                'sub_section': sub_section,
                'source': '신한라이프 2023 ESG 보고서'
            }
        }
        chunks.append(chunk)
    
    return chunks

def extract_esg_performance_base(ppt_path):
    """PPT 파일에서 ESG 관련 내용을 추출합니다."""
    print(f"PPT 파일을 읽는 중: {ppt_path}")
    
    prs = Presentation(ppt_path)
    num_slides = len(prs.slides)
    print(f"PPT 총 슬라이드 수: {num_slides}")
    
    # 섹션 범위 정의 (슬라이드 번호는 0부터 시작)
    section_ranges = {
        'Environmental': (24, 38),  # 25~39
        'Social': (40, 79),         # 41~80
        'Governance': (81, 99),     # 82~100
    }
    
    all_chunks = []
    
    for section_type, (start, end) in section_ranges.items():
        for slide_idx in range(start, end + 1):
            if slide_idx < num_slides:
                print(f"{section_type} 섹션 - 슬라이드 {slide_idx + 1} 처리 중...")
                slide = prs.slides[slide_idx]
                slide_text = extract_slide_content(slide)
                
                if slide_text.strip():
                    # 슬라이드 제목을 sub_section으로 사용
                    sub_section = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_idx + 1}"
                    chunks = create_chunks_base(slide_text, section_type, sub_section)
                    all_chunks.extend(chunks)
                    print(f"{section_type} 섹션({sub_section})에서 {len(chunks)}개의 청크 생성됨")
    
    return all_chunks

def main():
    ppt_path = "temp/data/ppts/2023 신한라이프 지속가능경영보고서.pptx"
    print("프로그램 시작")
    
    # PPT 파일이 존재하는지 확인
    if not Path(ppt_path).exists():
        print(f"오류: PPT 파일을 찾을 수 없습니다: {ppt_path}")
        return
    
    chunks = extract_esg_performance_base(ppt_path)
    
    if chunks:
        output_path = "temp/data/processed/esg_chunks_ppt_base.json"
        # 출력 디렉토리가 없으면 생성
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(chunks, f, ensure_ascii=False, indent=2)
        print(f"ESG 데이터가 {len(chunks)}개의 청크로 나뉘어 {output_path}에 저장되었습니다.")
    else:
        print("ESG Performance 섹션을 찾을 수 없습니다.")

if __name__ == "__main__":
    main() 